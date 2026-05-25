from __future__ import annotations

from pathlib import Path
from typing import Any
import json
import re
import time

from openai import OpenAI

from pypdf import PdfReader

from docx import Document
from docx.document import Document as DocxDocument
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph

from pptx import Presentation
from pptx.table import Table as PptxTable

from openpyxl import load_workbook
import xlrd


# ============================================================
# CONFIG
# ============================================================

OPENAI_API_KEY = "sk-PASTE_YOUR_KEY_HERE"

INPUT_DIR = Path(r"C:\docs")
OUTPUT_DIR = Path(r"C:\rag_json")

MODEL = "gpt-4.1-mini"

MAX_CHARS_PER_PART = 30_000
REQUEST_SLEEP_SECONDS = 0.2

SAVE_EXTRACTED_TEXT = True
OVERWRITE_EXISTING_JSON = True

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".xls",
}


# ============================================================
# PROMPT + SCHEMA
# ============================================================

PROMPT = """
Ты парсер документов для RAG-базы.

Я передам тебе текст, извлечённый из документа.

В тексте могут быть технические маркеры:
[PAGE 1]
[PARAGRAPH 12]
[STYLE Heading 1]
[LIST]
[SLIDE 3]
[SHEET Sheet1]
[ROW 15]
[TABLE 2]

Твоя задача — разбить документ на логические блоки.

Правила:
1. Не пересказывай текст.
2. Не добавляй ничего от себя.
3. Не исправляй смысл.
4. Сохраняй текст максимально дословно.
5. Не добавляй технические маркеры в поле text.
6. Если видишь [LIST], почти всегда ставь type = "list_item".
7. Не делай отдельный блок на каждую новую строку.
8. Объединяй соседние строки, если они относятся к одному абзацу, таблице, слайду или смысловой секции.
9. Для Excel строки одной таблицы объединяй в type = "table", если они относятся к одной таблице.
10. type должен быть одним из:
    title, heading, paragraph, list_item, table, slide_text, unknown.
11. heading — текущий заголовок раздела или null.
12. page — номер страницы PDF или null.
13. paragraph — номер параграфа DOCX или null.
14. slide — номер слайда PPTX или null.
15. sheet — имя листа Excel или null.
16. row — номер начальной строки Excel или null.
17. Верни только JSON по схеме.
"""

JSON_OUTPUT_RULES = """
Верни только валидный JSON object.
Без markdown.
Без ```json.
Без комментариев.

Формат ответа:

{
  "document_title": "string",
  "document_type": "string",
  "language": "string",
  "blocks": [
    {
      "type": "title | heading | paragraph | list_item | table | slide_text | unknown",
      "heading": "string or null",
      "page": "integer or null",
      "paragraph": "integer or null",
      "slide": "integer or null",
      "sheet": "string or null",
      "row": "integer or null",
      "text": "string"
    }
  ]
}
"""

# ============================================================
# TEXT HELPERS
# ============================================================

def clean_text(value: Any) -> str:
    text = "" if value is None else str(value)
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def safe_filename(name: str) -> str:
    name = re.sub(r'[<>:"/\\|?*]', "_", name)
    name = name.strip()
    return name or "document"


def format_cell_value(value: Any) -> str:
    if value is None:
        return ""

    if isinstance(value, float) and value.is_integer():
        return str(int(value))

    return clean_text(value)


def is_supported_file(path: Path) -> bool:
    return path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS


# ============================================================
# PDF
# ============================================================

def read_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    parts: list[str] = []

    for page_number, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text() or "")

        if text:
            parts.append(f"[PAGE {page_number}]\n{text}")

    return "\n\n".join(parts)


# ============================================================
# DOCX
# ============================================================

def iter_docx_blocks(document: DocxDocument):
    body = document.element.body

    for child in body.iterchildren():
        if child.tag.endswith("}p"):
            yield Paragraph(child, document)
        elif child.tag.endswith("}tbl"):
            yield DocxTable(child, document)


def get_docx_style_name(paragraph: Paragraph) -> str:
    if paragraph.style is None:
        return "Unknown"

    return paragraph.style.name or "Unknown"


def is_docx_list_paragraph(paragraph: Paragraph) -> bool:
    p = paragraph._p
    p_pr = p.pPr

    if p_pr is not None and p_pr.numPr is not None:
        return True

    style_name = get_docx_style_name(paragraph).lower()

    list_words = [
        "list",
        "список",
        "перечень",
    ]

    return any(word in style_name for word in list_words)


def read_docx_table(table: DocxTable) -> str:
    rows: list[str] = []

    for row in table.rows:
        cells: list[str] = []

        for cell in row.cells:
            cells.append(clean_text(cell.text))

        row_text = " | ".join(cells).strip()

        if row_text:
            rows.append(row_text)

    return clean_text("\n".join(rows))


def read_docx(path: Path) -> str:
    document = Document(str(path))

    parts: list[str] = []
    paragraph_number = 0
    table_number = 0

    for block in iter_docx_blocks(document):
        if isinstance(block, Paragraph):
            text = clean_text(block.text)

            if text:
                style_name = get_docx_style_name(block)

                markers = [
                    f"[PARAGRAPH {paragraph_number}]",
                    f"[STYLE {style_name}]",
                ]

                if is_docx_list_paragraph(block):
                    markers.append("[LIST]")

                parts.append("\n".join(markers) + "\n" + text)

            paragraph_number += 1

        elif isinstance(block, DocxTable):
            table_text = read_docx_table(block)

            if table_text:
                parts.append(
                    f"[PARAGRAPH {paragraph_number}]\n"
                    f"[TABLE {table_number}]\n"
                    f"{table_text}"
                )

            paragraph_number += 1
            table_number += 1

    return "\n\n".join(parts)


# ============================================================
# PPTX
# ============================================================

def read_pptx_table(table: PptxTable) -> str:
    rows: list[str] = []

    for row in table.rows:
        cells: list[str] = []

        for cell in row.cells:
            cells.append(clean_text(cell.text))

        row_text = " | ".join(cells).strip()

        if row_text:
            rows.append(row_text)

    return clean_text("\n".join(rows))


def read_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []

        for shape_index, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_table", False):
                table_text = read_pptx_table(shape.table)

                if table_text:
                    slide_parts.append(
                        f"[SLIDE {slide_number}]\n"
                        f"[TABLE {shape_index}]\n"
                        f"{table_text}"
                    )

            elif getattr(shape, "has_text_frame", False):
                text = clean_text(shape.text)

                if text:
                    slide_parts.append(
                        f"[SLIDE {slide_number}]\n"
                        f"{text}"
                    )

        if slide_parts:
            parts.append("\n\n".join(slide_parts))

    return "\n\n".join(parts)


# ============================================================
# XLSX
# ============================================================

def read_xlsx(path: Path) -> str:
    workbook = load_workbook(
        filename=str(path),
        read_only=True,
        data_only=True,
    )

    parts: list[str] = []

    for sheet in workbook.worksheets:
        sheet_name = sheet.title

        for row in sheet.iter_rows():
            values = [format_cell_value(cell.value) for cell in row]

            while values and values[-1] == "":
                values.pop()

            if not any(values):
                continue

            row_number = row[0].row if row else None
            row_text = " | ".join(values)

            parts.append(
                f"[SHEET {sheet_name}]\n"
                f"[ROW {row_number}]\n"
                f"{row_text}"
            )

    workbook.close()

    return "\n\n".join(parts)


# ============================================================
# XLS
# ============================================================

def read_xls(path: Path) -> str:
    workbook = xlrd.open_workbook(str(path))
    parts: list[str] = []

    for sheet in workbook.sheets():
        sheet_name = sheet.name

        for row_index in range(sheet.nrows):
            values: list[str] = []

            for col_index in range(sheet.ncols):
                value = sheet.cell_value(row_index, col_index)
                values.append(format_cell_value(value))

            while values and values[-1] == "":
                values.pop()

            if not any(values):
                continue

            row_number = row_index + 1
            row_text = " | ".join(values)

            parts.append(
                f"[SHEET {sheet_name}]\n"
                f"[ROW {row_number}]\n"
                f"{row_text}"
            )

    return "\n\n".join(parts)


# ============================================================
# READ DOCUMENT
# ============================================================

def read_document(path: Path) -> str:
    extension = path.suffix.lower()

    if extension == ".pdf":
        return read_pdf(path)

    if extension == ".docx":
        return read_docx(path)

    if extension == ".pptx":
        return read_pptx(path)

    if extension == ".xlsx":
        return read_xlsx(path)

    if extension == ".xls":
        return read_xls(path)

    raise ValueError(f"Unsupported file extension: {path.suffix}")


def find_documents(input_dir: Path) -> list[Path]:
    files = [path for path in input_dir.rglob("*") if is_supported_file(path)]
    return sorted(files)


# ============================================================
# SPLIT
# ============================================================

def split_text(text: str, max_chars: int) -> list[str]:
    raw_blocks = text.split("\n\n")

    parts: list[str] = []
    current = ""

    for block in raw_blocks:
        block = block.strip()

        if not block:
            continue

        candidate = f"{current}\n\n{block}" if current else block

        if len(candidate) <= max_chars:
            current = candidate
        else:
            if current:
                parts.append(current)

            current = block

    if current:
        parts.append(current)

    return parts


# ============================================================
# OPENAI
# ============================================================

def create_client() -> OpenAI:
    if not OPENAI_API_KEY or OPENAI_API_KEY == "sk-PASTE_YOUR_KEY_HERE":
        raise RuntimeError("Set OPENAI_API_KEY at the top of the script.")

    return OpenAI(api_key=OPENAI_API_KEY)


def parse_part_with_llm(
    client: OpenAI,
    file_name: str,
    text_part: str,
    part_index: int,
    total_parts: int,
) -> dict[str, Any]:
    system_prompt = PROMPT + "\n\n" + JSON_OUTPUT_RULES

    user_prompt = f"""
    FILE_NAME: {file_name}
    PART: {part_index}/{total_parts}
    
    TEXT:
    {text_part}
    """

    response = client.chat.completions.create(
        model=MODEL,
        messages=[
            {
                "role": "system",
                "content": system_prompt,
            },
            {
                "role": "user",
                "content": user_prompt,
            },
        ],
        response_format={
            "type": "json_object",
        },
    )

    content = response.choices[0].message.content

    if content is None:
        raise RuntimeError("LLM returned empty response")

    return json.loads(content)


# ============================================================
# PROCESSING
# ============================================================

def build_empty_document(path: Path) -> dict[str, Any]:
    return {
        "source_file": str(path),
        "source_file_name": path.name,
        "file_extension": path.suffix.lower(),
        "document_title": path.stem,
        "document_type": "unknown",
        "language": "unknown",
        "blocks": [],
    }


def add_block_metadata(
    block: dict[str, Any],
    source_path: Path,
    part_index: int,
    block_index: int,
) -> dict[str, Any]:
    block["block_id"] = f"{source_path.stem}::b{block_index:06d}"
    block["block_index"] = block_index
    block["source_file"] = str(source_path)
    block["source_file_name"] = source_path.name
    block["file_extension"] = source_path.suffix.lower()
    block["part"] = part_index
    return block


def save_extracted_text(path: Path, text: str) -> None:
    if not SAVE_EXTRACTED_TEXT:
        return

    extracted_dir = OUTPUT_DIR / "_extracted_text"
    extracted_dir.mkdir(parents=True, exist_ok=True)

    output_path = extracted_dir / f"{safe_filename(path.stem)}.txt"

    with open(output_path, "w", encoding="utf-8") as file:
        file.write(text)


def process_file(client: OpenAI, path: Path) -> dict[str, Any]:
    print(f"\nFILE: {path.name}")

    document = build_empty_document(path)

    raw_text = read_document(path)
    save_extracted_text(path, raw_text)

    if not raw_text:
        print("  EMPTY TEXT")
        return document

    text_parts = split_text(raw_text, MAX_CHARS_PER_PART)

    for part_index, text_part in enumerate(text_parts, start=1):
        print(f"  LLM PART {part_index}/{len(text_parts)}")

        parsed = parse_part_with_llm(
            client=client,
            file_name=path.name,
            text_part=text_part,
            part_index=part_index,
            total_parts=len(text_parts),
        )

        if part_index == 1:
            document["document_title"] = parsed.get("document_title") or path.stem
            document["document_type"] = parsed.get("document_type") or "unknown"
            document["language"] = parsed.get("language") or "unknown"

        for block in parsed.get("blocks", []):
            block_index = len(document["blocks"])

            block = add_block_metadata(
                block=block,
                source_path=path,
                part_index=part_index,
                block_index=block_index,
            )

            document["blocks"].append(block)

        if REQUEST_SLEEP_SECONDS > 0:
            time.sleep(REQUEST_SLEEP_SECONDS)

    return document


# ============================================================
# SAVE
# ============================================================

def save_document_json(path: Path, document: dict[str, Any]) -> None:
    output_name = f"{safe_filename(path.stem)}.json"
    output_path = OUTPUT_DIR / output_name

    if output_path.exists() and not OVERWRITE_EXISTING_JSON:
        print(f"  SKIP SAVE, EXISTS: {output_path}")
        return

    with open(output_path, "w", encoding="utf-8") as file:
        json.dump(document, file, ensure_ascii=False, indent=2)

    print(f"  SAVED JSON: {output_path}")


def save_all_blocks(documents: list[dict[str, Any]]) -> None:
    output_path = OUTPUT_DIR / "all_blocks.jsonl"

    with open(output_path, "w", encoding="utf-8") as file:
        for document in documents:
            for block in document["blocks"]:
                file.write(json.dumps(block, ensure_ascii=False) + "\n")

    print(f"\nSAVED BLOCKS: {output_path}")


# ============================================================
# MAIN
# ============================================================

def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    client = create_client()

    files = find_documents(INPUT_DIR)

    print(f"INPUT_DIR: {INPUT_DIR}")
    print(f"OUTPUT_DIR: {OUTPUT_DIR}")
    print(f"FOUND FILES: {len(files)}")

    documents: list[dict[str, Any]] = []

    for file_path in files:
        try:
            document = process_file(client, file_path)
            save_document_json(file_path, document)
            documents.append(document)

        except Exception as error:
            print(f"\nERROR: {file_path}")
            print(type(error).__name__)
            print(error)

    save_all_blocks(documents)

    print("\nDONE")


if __name__ == "__main__":
    main()