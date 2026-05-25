from __future__ import annotations

from pathlib import Path
from typing import Any
import hashlib
import json
import re

from pypdf import PdfReader

from docx import Document
from docx.document import Document as DocxDocument
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph

from pptx import Presentation
from pptx.table import Table as PptxTable

from openpyxl import load_workbook
import xlrd


# ============================================================
# CONFIG
# ============================================================

INPUT_DIR = Path(r"C:\docs")
OUTPUT_DIR = Path(r"C:\rag_json")

SAVE_EXTRACTED_TEXT = True
OVERWRITE_EXISTING_JSON = True

MAX_EXCEL_ROWS_PER_BLOCK = 50

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".xls",
}


# ============================================================
# HELPERS
# ============================================================

def clean_text(value: Any) -> str:
    text = "" if value is None else str(value)
    text = text.replace("\x00", "")
    text = text.replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def safe_filename(name: str) -> str:
    name = re.sub(r'[<>:"/\\|?*]', "_", name)
    name = name.strip()
    return name or "document"


def stable_id(value: str) -> str:
    return hashlib.sha256(value.encode("utf-8")).hexdigest()[:16]


def make_document_id(path: Path) -> str:
    try:
        relative = path.relative_to(INPUT_DIR)
        value = str(relative).lower()
    except ValueError:
        value = str(path).lower()

    return stable_id(value)


def is_supported_file(path: Path) -> bool:
    return path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS


def find_documents(input_dir: Path) -> list[Path]:
    files = [path for path in input_dir.rglob("*") if is_supported_file(path)]
    return sorted(files)


def format_cell_value(value: Any) -> str:
    if value is None:
        return ""

    if isinstance(value, float) and value.is_integer():
        return str(int(value))

    return clean_text(value)


def build_block(
    document_id: str,
    path: Path,
    block_index: int,
    block_type: str,
    text: str,
    heading: str | None = None,
    heading_path: list[str] | None = None,
    page: int | None = None,
    paragraph: int | None = None,
    slide: int | None = None,
    sheet: str | None = None,
    row_start: int | None = None,
    row_end: int | None = None,
    extra: dict[str, Any] | None = None,
) -> dict[str, Any]:
    block = {
        "block_id": f"{document_id}::b{block_index:06d}",
        "document_id": document_id,
        "source_file": str(path),
        "source_file_name": path.name,
        "file_extension": path.suffix.lower(),
        "block_index": block_index,
        "type": block_type,
        "heading": heading,
        "heading_path": heading_path or [],
        "page": page,
        "paragraph": paragraph,
        "slide": slide,
        "sheet": sheet,
        "row_start": row_start,
        "row_end": row_end,
        "text": clean_text(text),
    }

    if extra:
        block.update(extra)

    return block


def is_list_text(text: str) -> bool:
    patterns = [
        r"^[-–—•●▪]\s+",
        r"^\d+[\.\)]\s+",
        r"^\d+\.\d+[\.\)]?\s+",
        r"^[a-zA-Zа-яА-Я]\)\s+",
    ]

    return any(re.match(pattern, text.strip()) for pattern in patterns)


# ============================================================
# PDF
# ============================================================

def guess_pdf_block_type(text: str) -> str:
    stripped = text.strip()

    if not stripped:
        return "unknown"

    if is_list_text(stripped):
        return "list_item"

    lines = [line.strip() for line in stripped.splitlines() if line.strip()]

    if len(lines) == 1:
        line = lines[0]

        if len(line) <= 100 and not line.endswith((".", ",", ";")):
            return "heading"

    return "paragraph"


def split_pdf_page_text(text: str) -> list[str]:
    text = clean_text(text)

    if not text:
        return []

    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

    if len(paragraphs) > 1:
        return paragraphs

    lines = [line.strip() for line in text.splitlines() if line.strip()]

    if not lines:
        return []

    blocks: list[str] = []
    current: list[str] = []

    for line in lines:
        current.append(line)

        if line.endswith((".", "!", "?", ":", ";")):
            blocks.append(" ".join(current))
            current = []

    if current:
        blocks.append(" ".join(current))

    return blocks


def parse_pdf(path: Path, document_id: str) -> list[dict[str, Any]]:
    reader = PdfReader(str(path))
    blocks: list[dict[str, Any]] = []

    current_heading: str | None = None
    heading_path: list[str] = []

    for page_number, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        page_blocks = split_pdf_page_text(page_text)

        for text in page_blocks:
            block_type = guess_pdf_block_type(text)

            if block_type == "heading":
                current_heading = text
                heading_path = [text]

            block = build_block(
                document_id=document_id,
                path=path,
                block_index=len(blocks),
                block_type=block_type,
                text=text,
                heading=current_heading,
                heading_path=heading_path,
                page=page_number,
            )

            blocks.append(block)

    return blocks


# ============================================================
# DOCX
# ============================================================

def iter_docx_blocks(document: DocxDocument):
    body = document.element.body

    for child in body.iterchildren():
        if child.tag.endswith("}p"):
            yield Paragraph(child, document)
        elif child.tag.endswith("}tbl"):
            yield DocxTable(child, document)


def get_docx_style_name(paragraph: Paragraph) -> str:
    if paragraph.style is None:
        return "Unknown"

    return paragraph.style.name or "Unknown"


def get_heading_level_from_style(style_name: str) -> int | None:
    style = style_name.lower().strip()

    match = re.search(r"heading\s*(\d+)", style)
    if match:
        return int(match.group(1))

    match = re.search(r"заголовок\s*(\d+)", style)
    if match:
        return int(match.group(1))

    if style in {"title", "название"}:
        return 0

    return None


def is_docx_list_paragraph(paragraph: Paragraph) -> bool:
    p = paragraph._p
    p_pr = p.pPr

    if p_pr is not None and p_pr.numPr is not None:
        return True

    style_name = get_docx_style_name(paragraph).lower()

    list_words = [
        "list",
        "список",
        "перечень",
    ]

    return any(word in style_name for word in list_words)


def read_docx_table(table: DocxTable) -> str:
    rows: list[str] = []

    for row in table.rows:
        cells: list[str] = []

        for cell in row.cells:
            cells.append(clean_text(cell.text))

        row_text = " | ".join(cells).strip()

        if row_text:
            rows.append(row_text)

    return clean_text("\n".join(rows))


def parse_docx(path: Path, document_id: str) -> list[dict[str, Any]]:
    document = Document(str(path))
    blocks: list[dict[str, Any]] = []

    heading_stack: dict[int, str] = {}
    current_heading: str | None = None

    paragraph_number = 0
    table_number = 0

    for item in iter_docx_blocks(document):
        if isinstance(item, Paragraph):
            text = clean_text(item.text)
            style_name = get_docx_style_name(item)
            heading_level = get_heading_level_from_style(style_name)

            if not text:
                paragraph_number += 1
                continue

            if heading_level is not None:
                if heading_level == 0:
                    block_type = "title"
                    heading_stack = {0: text}
                else:
                    block_type = "heading"
                    heading_stack[heading_level] = text

                    for level in list(heading_stack.keys()):
                        if level > heading_level:
                            del heading_stack[level]

                current_heading = text

            elif is_docx_list_paragraph(item) or is_list_text(text):
                block_type = "list_item"
            else:
                block_type = "paragraph"

            heading_path = [
                heading_stack[level]
                for level in sorted(heading_stack.keys())
                if level != 0 or block_type != "title"
            ]

            if block_type in {"title", "heading"}:
                block_heading = text
            else:
                block_heading = current_heading

            block = build_block(
                document_id=document_id,
                path=path,
                block_index=len(blocks),
                block_type=block_type,
                text=text,
                heading=block_heading,
                heading_path=heading_path,
                paragraph=paragraph_number,
                extra={
                    "style": style_name,
                    "is_list": block_type == "list_item",
                },
            )

            blocks.append(block)
            paragraph_number += 1

        elif isinstance(item, DocxTable):
            table_text = read_docx_table(item)

            if table_text:
                heading_path = [
                    heading_stack[level]
                    for level in sorted(heading_stack.keys())
                ]

                block = build_block(
                    document_id=document_id,
                    path=path,
                    block_index=len(blocks),
                    block_type="table",
                    text=table_text,
                    heading=current_heading,
                    heading_path=heading_path,
                    paragraph=paragraph_number,
                    extra={
                        "table_index": table_number,
                    },
                )

                blocks.append(block)

            paragraph_number += 1
            table_number += 1

    return blocks


# ============================================================
# PPTX
# ============================================================

def read_pptx_table(table: PptxTable) -> str:
    rows: list[str] = []

    for row in table.rows:
        cells: list[str] = []

        for cell in row.cells:
            cells.append(clean_text(cell.text))

        row_text = " | ".join(cells).strip()

        if row_text:
            rows.append(row_text)

    return clean_text("\n".join(rows))


def parse_pptx(path: Path, document_id: str) -> list[dict[str, Any]]:
    presentation = Presentation(str(path))
    blocks: list[dict[str, Any]] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text_parts: list[str] = []

        for shape_index, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_table", False):
                table_text = read_pptx_table(shape.table)

                if table_text:
                    block = build_block(
                        document_id=document_id,
                        path=path,
                        block_index=len(blocks),
                        block_type="table",
                        text=table_text,
                        slide=slide_number,
                        extra={
                            "shape_index": shape_index,
                        },
                    )

                    blocks.append(block)

            elif getattr(shape, "has_text_frame", False):
                text = clean_text(shape.text)

                if text:
                    slide_text_parts.append(text)

        slide_text = clean_text("\n".join(slide_text_parts))

        if slide_text:
            first_line = slide_text.splitlines()[0].strip()
            heading = first_line if len(first_line) <= 120 else None

            block = build_block(
                document_id=document_id,
                path=path,
                block_index=len(blocks),
                block_type="slide_text",
                text=slide_text,
                heading=heading,
                heading_path=[heading] if heading else [],
                slide=slide_number,
            )

            blocks.append(block)

    return blocks


# ============================================================
# EXCEL COMMON
# ============================================================

def rows_to_table_blocks(
    path: Path,
    document_id: str,
    sheet_name: str,
    rows: list[tuple[int, list[str]]],
    existing_blocks: list[dict[str, Any]],
) -> None:
    group: list[tuple[int, list[str]]] = []

    def flush_group() -> None:
        nonlocal group

        if not group:
            return

        for start in range(0, len(group), MAX_EXCEL_ROWS_PER_BLOCK):
            chunk = group[start:start + MAX_EXCEL_ROWS_PER_BLOCK]

            row_start = chunk[0][0]
            row_end = chunk[-1][0]

            lines = []

            for row_number, values in chunk:
                line = f"row {row_number}: " + " | ".join(values)
                lines.append(line)

            text = "\n".join(lines)

            block = build_block(
                document_id=document_id,
                path=path,
                block_index=len(existing_blocks),
                block_type="table",
                text=text,
                sheet=sheet_name,
                row_start=row_start,
                row_end=row_end,
                extra={
                    "excel_table_rows": len(chunk),
                },
            )

            existing_blocks.append(block)

        group = []

    for row_number, values in rows:
        is_empty = not any(value.strip() for value in values)

        if is_empty:
            flush_group()
            continue

        while values and values[-1] == "":
            values.pop()

        group.append((row_number, values))

    flush_group()


# ============================================================
# XLSX
# ============================================================

def parse_xlsx(path: Path, document_id: str) -> list[dict[str, Any]]:
    workbook = load_workbook(
        filename=str(path),
        read_only=True,
        data_only=True,
    )

    blocks: list[dict[str, Any]] = []

    for sheet in workbook.worksheets:
        rows: list[tuple[int, list[str]]] = []

        for row in sheet.iter_rows():
            row_number = row[0].row if row else 0
            values = [format_cell_value(cell.value) for cell in row]
            rows.append((row_number, values))

        rows_to_table_blocks(
            path=path,
            document_id=document_id,
            sheet_name=sheet.title,
            rows=rows,
            existing_blocks=blocks,
        )

    workbook.close()
    return blocks


# ============================================================
# XLS
# ============================================================

def parse_xls(path: Path, document_id: str) -> list[dict[str, Any]]:
    workbook = xlrd.open_workbook(str(path))
    blocks: list[dict[str, Any]] = []

    for sheet in workbook.sheets():
        rows: list[tuple[int, list[str]]] = []

        for row_index in range(sheet.nrows):
            row_number = row_index + 1
            values: list[str] = []

            for col_index in range(sheet.ncols):
                value = sheet.cell_value(row_index, col_index)
                values.append(format_cell_value(value))

            rows.append((row_number, values))

        rows_to_table_blocks(
            path=path,
            document_id=document_id,
            sheet_name=sheet.name,
            rows=rows,
            existing_blocks=blocks,
        )

    return blocks


# ============================================================
# DOCUMENT PARSER
# ============================================================

def parse_document(path: Path) -> dict[str, Any]:
    document_id = make_document_id(path)
    extension = path.suffix.lower()

    if extension == ".pdf":
        blocks = parse_pdf(path, document_id)
    elif extension == ".docx":
        blocks = parse_docx(path, document_id)
    elif extension == ".pptx":
        blocks = parse_pptx(path, document_id)
    elif extension == ".xlsx":
        blocks = parse_xlsx(path, document_id)
    elif extension == ".xls":
        blocks = parse_xls(path, document_id)
    else:
        raise ValueError(f"Unsupported file extension: {extension}")

    return {
        "document_id": document_id,
        "document_title": path.stem,
        "source_file": str(path),
        "source_file_name": path.name,
        "file_extension": extension,
        "block_count": len(blocks),
        "blocks": blocks,
    }


# ============================================================
# SAVE
# ============================================================

def save_document_json(path: Path, document: dict[str, Any]) -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    output_name = f"{safe_filename(path.stem)}.json"
    output_path = OUTPUT_DIR / output_name

    if output_path.exists() and not OVERWRITE_EXISTING_JSON:
        print(f"  SKIP SAVE, EXISTS: {output_path}")
        return

    with open(output_path, "w", encoding="utf-8") as file:
        json.dump(document, file, ensure_ascii=False, indent=2)

    print(f"  SAVED JSON: {output_path}")


def save_all_blocks(documents: list[dict[str, Any]]) -> None:
    output_path = OUTPUT_DIR / "all_blocks.jsonl"

    with open(output_path, "w", encoding="utf-8") as file:
        for document in documents:
            for block in document["blocks"]:
                file.write(json.dumps(block, ensure_ascii=False) + "\n")

    print(f"\nSAVED BLOCKS: {output_path}")


def save_extracted_text(document: dict[str, Any]) -> None:
    if not SAVE_EXTRACTED_TEXT:
        return

    extracted_dir = OUTPUT_DIR / "_extracted_text"
    extracted_dir.mkdir(parents=True, exist_ok=True)

    output_path = extracted_dir / f"{safe_filename(document['document_title'])}.txt"

    lines: list[str] = []

    for block in document["blocks"]:
        location_parts = []

        if block.get("page") is not None:
            location_parts.append(f"page={block['page']}")

        if block.get("paragraph") is not None:
            location_parts.append(f"paragraph={block['paragraph']}")

        if block.get("slide") is not None:
            location_parts.append(f"slide={block['slide']}")

        if block.get("sheet") is not None:
            location_parts.append(f"sheet={block['sheet']}")

        if block.get("row_start") is not None:
            location_parts.append(f"rows={block['row_start']}-{block['row_end']}")

        location = ", ".join(location_parts)

        lines.append(
            f"[{block['block_id']}] "
            f"type={block['type']} "
            f"{location}\n"
            f"{block['text']}\n"
        )

    with open(output_path, "w", encoding="utf-8") as file:
        file.write("\n\n".join(lines))


# ============================================================
# MAIN
# ============================================================

def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    files = find_documents(INPUT_DIR)

    print(f"INPUT_DIR: {INPUT_DIR}")
    print(f"OUTPUT_DIR: {OUTPUT_DIR}")
    print(f"FOUND FILES: {len(files)}")

    documents: list[dict[str, Any]] = []

    for file_path in files:
        print(f"\nFILE: {file_path.name}")

        try:
            document = parse_document(file_path)

            save_document_json(file_path, document)
            save_extracted_text(document)

            documents.append(document)

            print(f"  BLOCKS: {document['block_count']}")

        except Exception as error:
            print(f"  ERROR: {type(error).__name__}")
            print(f"  {error}")

    save_all_blocks(documents)

    print("\nDONE")


if __name__ == "__main__":
    main()